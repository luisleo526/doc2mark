"""Office image-dominance route: probe classification + gating + graceful fallback."""
import io
from pathlib import Path
from unittest.mock import patch

import pytest

from doc2mark.formats.office import OfficeProcessor

SAMP = Path("sample_documents")


class _StubOCR:
    config = None


@pytest.fixture
def image_pptx(tmp_path):
    """A 1-slide pptx that is a single full-bleed image -> image-dominant."""
    pptx = pytest.importorskip("pptx")
    Image = pytest.importorskip("PIL.Image")
    from pptx.util import Emu
    prs = pptx.Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    Image.new("RGB", (800, 600), "navy").save(buf, format="PNG")
    s.shapes.add_picture(io.BytesIO(buf.getvalue()), 0, 0,
                         width=prs.slide_width, height=prs.slide_height)
    out = tmp_path / "img.pptx"
    prs.save(str(out))
    return out


def test_probe_image_dominant_vs_text(image_pptx):
    p = OfficeProcessor(ocr=_StubOCR())
    assert p._is_image_dominant(image_pptx) is True
    assert p._is_image_dominant(SAMP / "sample_document.docx") is False
    assert p._is_image_dominant(SAMP / "sample_presentation.pptx") is False


def test_xlsx_never_routes():
    p = OfficeProcessor(ocr=_StubOCR())
    assert p._maybe_route_image_dominant(
        SAMP / "sample_spreadsheet.xlsx", 100, ocr_images=True, extract_images=True) is None


def test_no_route_without_ocr(image_pptx):
    assert OfficeProcessor(ocr=None)._maybe_route_image_dominant(
        image_pptx, 100, ocr_images=True, extract_images=True) is None


def test_no_route_when_ocr_not_requested(image_pptx):
    p = OfficeProcessor(ocr=_StubOCR())
    assert p._maybe_route_image_dominant(
        image_pptx, 100, ocr_images=False, extract_images=True) is None
    assert p._maybe_route_image_dominant(
        image_pptx, 100, ocr_images=True, extract_images=False) is None


def test_text_doc_not_routed():
    p = OfficeProcessor(ocr=_StubOCR())
    assert p._maybe_route_image_dominant(
        SAMP / "sample_document.docx", 100, ocr_images=True, extract_images=True) is None


def test_route_falls_back_on_conversion_failure(image_pptx):
    """Image-dominant + OCR, but conversion raises -> None (native fallback), never raises."""
    p = OfficeProcessor(ocr=_StubOCR())
    with patch.object(OfficeProcessor, "_process_as_image_dominant",
                      side_effect=RuntimeError("no soffice")):
        assert p._maybe_route_image_dominant(
            image_pptx, 100, ocr_images=True, extract_images=True) is None
