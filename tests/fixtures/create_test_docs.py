#!/usr/bin/env python3
"""Generate test documents with known content for end-to-end RAG feature tests.

Run this script once to create the fixture documents under tests/fixtures/.
The generated files are committed to the repo so tests don't need runtime generation.
"""

import os
import sys
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


# ---------------------------------------------------------------------------
# 1. DOCX with footnotes, endnotes, page breaks, headers, footers, sections
# ---------------------------------------------------------------------------

def create_docx_with_footnotes():
    """Create a DOCX that exercises all RAG features."""
    import zipfile
    from docx import Document
    from docx.oxml.ns import qn
    from lxml import etree

    doc = Document()

    # --- Header / Footer ---
    section = doc.sections[0]
    header = section.header
    header.paragraphs[0].text = "ACME Corp — Confidential"
    footer = section.footer
    footer.paragraphs[0].text = "Page Footer Text"

    # --- Page 1 content ---
    doc.add_heading("Introduction", level=1)
    p1 = doc.add_paragraph("This is the first paragraph on page one.")
    # We'll inject a footnote reference into p1 via raw XML below.
    doc.add_paragraph("Some more text on the first page with details.")

    # --- Explicit page break ---
    doc.add_page_break()

    # --- Page 2 content ---
    doc.add_heading("Methodology", level=1)
    p2 = doc.add_paragraph("The methodology section begins on page two.")
    doc.add_paragraph("We use a novel approach combining X and Y.")

    # --- Another page break ---
    doc.add_page_break()

    # --- Page 3 content ---
    doc.add_heading("Results", level=1)
    doc.add_paragraph("Results are presented in the table below.")

    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    cells = table.rows[0].cells
    cells[0].text = "Metric"
    cells[1].text = "Value"
    cells = table.rows[1].cells
    cells[0].text = "Accuracy"
    cells[1].text = "95.2%"
    cells = table.rows[2].cells
    cells[0].text = "F1 Score"
    cells[1].text = "0.93"

    doc.add_heading("Conclusion", level=1)
    doc.add_paragraph("In conclusion, the approach works well.")

    # --- Save first, then inject footnotes via raw XML manipulation ---
    docx_path = FIXTURES_DIR / "test_rag_features.docx"
    doc.save(str(docx_path))

    # Now inject footnotes into the ZIP
    _inject_footnotes_into_docx(docx_path, p1_text="first paragraph on page one")

    print(f"  Created: {docx_path}")
    return docx_path


def _inject_footnotes_into_docx(docx_path, p1_text="first paragraph"):
    """Inject footnotes.xml and endnotes.xml + references into the saved DOCX."""
    import zipfile
    import shutil
    from lxml import etree

    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    CT = "http://schemas.openxmlformats.org/package/2006/content-types"

    # Read existing ZIP
    tmp_path = docx_path.with_suffix(".tmp")
    shutil.copy(docx_path, tmp_path)

    with zipfile.ZipFile(tmp_path, "r") as zin, zipfile.ZipFile(docx_path, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == "[Content_Types].xml":
                tree = etree.fromstring(data)
                # Add content type for footnotes and endnotes
                etree.SubElement(tree, "Override",
                    PartName="/word/footnotes.xml",
                    ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.footnotes+xml")
                etree.SubElement(tree, "Override",
                    PartName="/word/endnotes.xml",
                    ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.endnotes+xml")
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)

            elif item.filename == "word/_rels/document.xml.rels":
                tree = etree.fromstring(data)
                ns = {"": R}
                # Find max rId
                max_id = 0
                for rel in tree:
                    rid = rel.get("Id", "")
                    if rid.startswith("rId"):
                        try:
                            max_id = max(max_id, int(rid[3:]))
                        except ValueError:
                            pass
                etree.SubElement(tree, "Relationship",
                    Id=f"rId{max_id+1}",
                    Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes",
                    Target="footnotes.xml")
                etree.SubElement(tree, "Relationship",
                    Id=f"rId{max_id+2}",
                    Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/endnotes",
                    Target="endnotes.xml")
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)

            elif item.filename == "word/document.xml":
                tree = etree.fromstring(data)
                # Find the first paragraph that contains our target text and
                # inject a footnote reference run
                for p_el in tree.iter(f"{{{W}}}p"):
                    full_text = "".join(t.text or "" for t in p_el.iter(f"{{{W}}}t"))
                    if p1_text in full_text:
                        # Add a new run with footnoteReference
                        new_run = etree.SubElement(p_el, f"{{{W}}}r")
                        rpr = etree.SubElement(new_run, f"{{{W}}}rPr")
                        etree.SubElement(rpr, f"{{{W}}}rStyle", {f"{{{W}}}val": "FootnoteReference"})
                        etree.SubElement(new_run, f"{{{W}}}footnoteReference", {f"{{{W}}}id": "1"})
                        break
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)

            zout.writestr(item, data)

        # Add footnotes.xml
        fn_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:footnotes xmlns:w="{W}">'
            f'<w:footnote w:type="separator" w:id="0">'
            f'<w:p><w:r><w:separator/></w:r></w:p>'
            f'</w:footnote>'
            f'<w:footnote w:type="continuationSeparator" w:id="-1">'
            f'<w:p><w:r><w:continuationSeparator/></w:r></w:p>'
            f'</w:footnote>'
            f'<w:footnote w:id="1">'
            f'<w:p><w:r><w:t>This is a test footnote about the first paragraph.</w:t></w:r></w:p>'
            f'</w:footnote>'
            f'</w:footnotes>'
        )
        zout.writestr("word/footnotes.xml", fn_xml)

        # Add endnotes.xml
        en_xml = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:endnotes xmlns:w="{W}">'
            f'<w:endnote w:type="separator" w:id="0">'
            f'<w:p><w:r><w:separator/></w:r></w:p>'
            f'</w:endnote>'
            f'<w:endnote w:type="continuationSeparator" w:id="-1">'
            f'<w:p><w:r><w:continuationSeparator/></w:r></w:p>'
            f'</w:endnote>'
            f'<w:endnote w:id="1">'
            f'<w:p><w:r><w:t>An endnote providing additional context.</w:t></w:r></w:p>'
            f'</w:endnote>'
            f'</w:endnotes>'
        )
        zout.writestr("word/endnotes.xml", en_xml)

    os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# 2. PDF with repeated headers, footers, footnotes, and multiple sections
# ---------------------------------------------------------------------------

def create_pdf_with_headers_and_footnotes():
    """Create a multi-page PDF with repeated headers/footers and footnotes."""
    import pymupdf

    pdf_path = FIXTURES_DIR / "test_rag_features.pdf"
    doc = pymupdf.open()

    header_text = "ACME Corp — Internal Report"
    footer_text = "Confidential — Do Not Distribute"

    pages_content = [
        {
            "title": "Executive Summary",
            "body": [
                "This report presents the findings of our Q3 analysis.",
                "Revenue grew by 15% compared to the previous quarter.",
                "Key drivers include expansion into new markets.",
            ],
            "footnote": "1. Source: Internal financial data, Q3 2025.",
        },
        {
            "title": "Market Analysis",
            "body": [
                "The total addressable market expanded to $4.2B in 2025.",
                "Our market share increased from 12% to 15%.",
                "Competitive pressure remains high in the enterprise segment.",
                "Customer acquisition cost decreased by 8%.",
            ],
            "footnote": "2. Based on Gartner market sizing report.",
        },
        {
            "title": "Product Development",
            "body": [
                "Three major features were shipped during Q3.",
                "Feature adoption rates exceeded 60% within the first month.",
                "Technical debt was reduced by 20% through refactoring.",
            ],
            "footnote": None,
        },
        {
            "title": "Financial Results",
            "body": [
                "Total revenue: $12.5M (up 15% QoQ).",
                "Operating margin: 22% (up from 19%).",
                "Cash position: $45M.",
            ],
            "footnote": "3. All figures are unaudited.",
        },
        {
            "title": "Outlook",
            "body": [
                "We expect continued growth in Q4 driven by seasonal demand.",
                "New product launches are planned for November.",
                "Hiring will focus on engineering and sales roles.",
            ],
            "footnote": None,
        },
    ]

    for page_data in pages_content:
        page = doc.new_page(width=612, height=792)  # US Letter

        # Header (top of every page)
        page.insert_text((72, 36), header_text, fontsize=9, color=(0.4, 0.4, 0.4))

        # Page title
        y = 72
        page.insert_text((72, y), page_data["title"], fontsize=18, color=(0, 0, 0))
        y += 36

        # Body paragraphs
        for para in page_data["body"]:
            page.insert_text((72, y), para, fontsize=11, color=(0, 0, 0))
            y += 20

        # Footnote at bottom (small text with numeric marker)
        if page_data["footnote"]:
            page.insert_text((72, 740), page_data["footnote"], fontsize=8, color=(0.3, 0.3, 0.3))

        # Footer (bottom of every page)
        page.insert_text((72, 770), footer_text, fontsize=8, color=(0.4, 0.4, 0.4))

    doc.save(str(pdf_path))
    doc.close()

    print(f"  Created: {pdf_path}")
    return pdf_path


# ---------------------------------------------------------------------------
# 3. XLSX with multiple sheets
# ---------------------------------------------------------------------------

def create_xlsx_with_sheets():
    """Create an XLSX with multiple named sheets."""
    import openpyxl

    xlsx_path = FIXTURES_DIR / "test_rag_features.xlsx"
    wb = openpyxl.Workbook()

    # Sheet 1: Sales Data
    ws1 = wb.active
    ws1.title = "Sales Data"
    ws1.append(["Product", "Q1", "Q2", "Q3"])
    ws1.append(["Widget A", 1000, 1200, 1500])
    ws1.append(["Widget B", 800, 950, 1100])
    ws1.append(["Widget C", 600, 700, 850])

    # Sheet 2: Expenses
    ws2 = wb.create_sheet("Expenses")
    ws2.append(["Category", "Amount"])
    ws2.append(["Salaries", 50000])
    ws2.append(["Rent", 10000])
    ws2.append(["Marketing", 15000])

    wb.save(str(xlsx_path))

    print(f"  Created: {xlsx_path}")
    return xlsx_path


# ---------------------------------------------------------------------------
# 4. Multi-section PPTX
# ---------------------------------------------------------------------------

def create_pptx_with_slides():
    """Create a PPTX with multiple slides and varied content."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pptx_path = FIXTURES_DIR / "test_rag_features.pptx"
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "Q3 2025 — Prepared by Analytics Team"

    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Overview"
    body = slide2.placeholders[1]
    body.text = "Total revenue reached $12.5M"
    body.text_frame.add_paragraph().text = "Growth rate: 15% QoQ"
    body.text_frame.add_paragraph().text = "Target exceeded by 3%"

    # Slide 3: Another content slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Metrics"
    body3 = slide3.placeholders[1]
    body3.text = "Customer count: 2,500"
    body3.text_frame.add_paragraph().text = "NPS: 72"
    body3.text_frame.add_paragraph().text = "Churn: 3.2%"

    prs.save(str(pptx_path))

    print(f"  Created: {pptx_path}")
    return pptx_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("Generating test fixture documents...")
    create_docx_with_footnotes()
    create_pdf_with_headers_and_footnotes()
    create_xlsx_with_sheets()
    create_pptx_with_slides()
    print("Done.")
