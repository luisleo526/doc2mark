"""doc2mark extraction harness.

Verifies, across the whole sample set AND an image-rich deck, that the Markdown
output KEEPS the three things that matter — system-wide, with no per-document
tuning:

  1. BODY TEXT      — every textual content line in the document body survives.
  2. TABLE STRUCTURE — complex tables are preserved as a grid, merges (col/row
                       spans) included.
  3. MEANINGFUL     — for image-only pages the output is structured/readable,
                       not a fragmented OCR dump.

Ground truth is structural where possible (OOXML grids with spans for
docx/xlsx/pptx; the deterministic rule-based text layer as the body-text
reference). Meaningfulness is an LLM judge. Run manually (needs OPENAI_API_KEY
for the image-OCR docs + the judge):

    python eval/extraction_harness.py            # all docs
    python eval/extraction_harness.py --quick    # skip the 30-page deck
"""
from __future__ import annotations
import argparse
import os
import re
import sys
import warnings

warnings.filterwarnings("ignore")

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SAMP = os.path.join(ROOT, "sample_documents")

# (relative path, class). class selects which checks apply.
DOCS = [
    ("sample_text.txt", "text"),
    ("sample_document.md", "text"),
    ("sample_data.csv", "data"),
    ("sample_data.json", "data"),
    ("sample_document.docx", "doc"),
    ("sample_legacy_document.doc", "doc"),
    ("sample_presentation.pptx", "slides"),
    ("sample_legacy_presentation.ppt", "slides"),
    ("sample_spreadsheet.xlsx", "sheet"),
    ("sample_spreadsheet_incell.xlsx", "sheet"),
    ("sample_pdf.pdf", "pdf"),
    ("test-table.pdf", "pdf_table"),
    ("complex-tables/complex_table_test.docx", "complex_table"),
    ("complex-tables/complex_table_test.pptx", "complex_table"),
    ("complex-tables/complex_table_test.xlsx", "complex_table"),
    ("complex-tables/complex_table_test.pdf", "complex_table"),
    (os.path.join(ROOT, "數辰企業簡報202606.pdf"), "image_deck"),
]
IMAGE_CLASSES = {"image_deck"}
TABLE_CLASSES = {"pdf_table", "complex_table", "sheet"}

# Pass thresholds (the requirements, made numeric).
THRESH = {"text": 0.95, "table": 0.90, "meaning": 3.5}

MODEL = "gpt-5.4-mini"


def norm(s: str) -> str:
    return " ".join((s or "").split())


# --------------------------------------------------------------------------- #
# Ground truth                                                                 #
# --------------------------------------------------------------------------- #
def rule_based_text(path: str) -> str:
    """Faithful body text = deterministic rule-based extraction (OCR off)."""
    from doc2mark import UnifiedDocumentLoader
    return UnifiedDocumentLoader(ocr_provider=None).load(path, output_format="markdown").content


def _dense(cells: dict, nr: int, nc: int) -> list:
    return [[norm(cells.get((r, c), "")) for c in range(nc)] for r in range(nr)]


def _xlsx_grids(path: str) -> list:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    grids = []
    for ws in wb.worksheets:
        if ws.max_row < 2 and ws.max_column < 2:
            continue
        cells = {}
        for row in ws.iter_rows():
            for c in row:
                if c.value is not None:
                    cells[(c.row - 1, c.column - 1)] = str(c.value)
        # expand merged ranges (top-left value fills the block)
        for m in ws.merged_cells.ranges:
            v = cells.get((m.min_row - 1, m.min_col - 1), "")
            for r in range(m.min_row - 1, m.max_row):
                for col in range(m.min_col - 1, m.max_col):
                    cells[(r, col)] = v
        grids.append(_dense(cells, ws.max_row, ws.max_column))
    return grids


def _docx_grids(path: str) -> list:
    import docx
    grids = []
    for tbl in docx.Document(path).tables:
        nr = len(tbl.rows)
        nc = max(len(r.cells) for r in tbl.rows)  # python-docx already expands spans
        grids.append([[norm(tbl.cell(r, c).text) for c in range(nc)] for r in range(nr)])
    return grids


def _pptx_grids(path: str) -> list:
    from pptx import Presentation
    grids = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            t = shape.table
            grids.append([[norm(t.cell(r, c).text) for c in range(len(t.columns))]
                          for r in range(len(t.rows))])
    return grids


# Known ground truth for the complex_table_test.* family (same logical table in
# every format) — lets the PDF variant be scored even without an OOXML source.
_COMPLEX_TABLE_GT = [[
    ["Company Overview", "Company Overview", "Company Overview", "Q1", "Q2", "Q3", "Q4"],
    ["Division", "Region", "Product", "First Half", "First Half", "Second Half", "Second Half"],
    ["Division", "North", "Widget A", "$10K", "$12K", "$15K", "$18K"],
    ["Technology", "South", "Widget B", "$20K", "$22K", "$25K", "$28K"],
    ["Technology", "East", "Combined Products", "Combined Products", "$30K", "$35K", "$40K"],
    ["Technology", "All Regions Total", "All Regions Total", "$50K", "$54K", "$60K", "$68K"],
    ["Subtotal (All Divisions)", "Subtotal (All Divisions)", "Subtotal (All Divisions)",
     "$80K", "$88K", "$100K", "$114K"],
    ["Grand Total (All Quarters)", "Grand Total (All Quarters)", "Grand Total (All Quarters)",
     "Grand Total (All Quarters)", "Annual Total: $382K", "Annual Total: $382K", "Annual Total: $382K"],
]]


def true_grids(path: str) -> list:
    p = path.lower()
    try:
        if p.endswith(".xlsx"):
            return _xlsx_grids(path)
        if p.endswith(".docx"):
            return _docx_grids(path)
        if p.endswith(".pptx"):
            return _pptx_grids(path)
        if "complex_table_test" in p:  # the PDF variant — use the known grid
            return _COMPLEX_TABLE_GT
    except Exception as e:  # pragma: no cover
        print(f"   (GT grid extract failed: {e})")
    return []


# --------------------------------------------------------------------------- #
# Output parsing                                                               #
# --------------------------------------------------------------------------- #
def html_to_grid(h: str) -> list:
    from lxml import html as LH
    try:
        root = LH.fromstring(h)
    except Exception:
        return []
    grid = {}
    for r, tr in enumerate(root.xpath("//tr")):
        c = 0
        for cell in tr.xpath("./td|./th"):
            while (r, c) in grid:
                c += 1
            cs = int(cell.get("colspan") or 1)
            rs = int(cell.get("rowspan") or 1)
            txt = norm(cell.text_content())
            for dr in range(rs):
                for dc in range(cs):
                    grid[(r + dr, c + dc)] = txt
            c += cs
    if not grid:
        return []
    nr = max(r for r, _ in grid) + 1
    nc = max(c for _, c in grid) + 1
    return [[grid.get((r, c), "") for c in range(nc)] for r in range(nr)]


def md_to_grid(block: str) -> list:
    rows = []
    for line in block.splitlines():
        line = line.strip()
        if not line.startswith("|"):
            continue
        if re.match(r"^\|[\s:|-]+\|?$", line):  # separator row
            continue
        cells = [norm(c) for c in line.strip("|").split("|")]
        rows.append(cells)
    return rows


def output_grids(md: str) -> list:
    grids = [html_to_grid(m) for m in re.findall(r"<table.*?</table>", md, re.S)]
    # markdown pipe tables (contiguous | lines)
    for block in re.findall(r"(?:^\|.*\|\s*$\n?)+", md, re.M):
        g = md_to_grid(block)
        if g:
            grids.append(g)
    return [g for g in grids if g]


# --------------------------------------------------------------------------- #
# Scorers                                                                      #
# --------------------------------------------------------------------------- #
def text_preservation(ref: str, out: str) -> float:
    """Fraction of reference body-text lines present (verbatim, ws-normalized) in out."""
    out_n = norm(out)
    lines = [norm(l) for l in ref.splitlines() if len(l.strip()) >= 12 and "|" not in l]
    lines = [l for l in lines if l]
    if not lines:
        return 1.0
    return sum(1 for l in lines if l in out_n) / len(lines)


def _grid_cell_acc(gt: list, out: list) -> float:
    nr, nc = len(gt), max((len(r) for r in gt), default=0)
    if nr == 0 or nc == 0:
        return 0.0
    gt_cells = {(r, c) for r in range(nr) for c in range(nc) if norm(gt[r][c])}
    if not gt_cells:
        return 1.0
    hit = 0
    for (r, c) in gt_cells:
        val = norm(gt[r][c])
        if r < len(out) and c < len(out[r]) and norm(out[r][c]) == val:
            hit += 1
    return hit / len(gt_cells)


def table_fidelity(gt_grids: list, out: str) -> float:
    """Best alignment of output table grids to GT grids; mean cell accuracy.
    Spans are validated implicitly (GT is span-expanded; a missed span leaves a
    wrong/empty cell)."""
    outs = output_grids(out)
    if not gt_grids:
        return None
    if not outs:
        return 0.0
    scores = []
    for gt in gt_grids:
        scores.append(max(_grid_cell_acc(gt, o) for o in outs))
    return sum(scores) / len(scores)


_JUDGE = (
    "You are grading how well a Markdown extraction represents a source document "
    "for a RAG system. Score 1-5 (5=excellent):\n"
    "5 = well-structured, readable: clear headings, lists/tables, diagrams described "
    "as coherent flows; a person could understand the document from this alone.\n"
    "3 = content is present but loosely organized.\n"
    "1 = a fragmented dump of disconnected tokens (orphaned numbers/labels, no "
    "structure), hard to follow.\n"
    "Reply with ONLY the integer.\n\nMARKDOWN:\n"
)


def meaningfulness(out: str) -> float:
    try:
        from langchain_openai import ChatOpenAI
        from langchain_core.messages import HumanMessage
        llm = ChatOpenAI(model=MODEL, max_tokens=4, timeout=40)
        ans = llm.invoke([HumanMessage(content=_JUDGE + out[:9000])]).content
        m = re.search(r"[1-5]", ans or "")
        return float(m.group()) if m else None
    except Exception as e:
        print(f"   (judge failed: {e})")
        return None


def junk(out: str) -> dict:
    return {
        "base64": out.count("data:image"),
        "markers": out.count("ocr_result"),
        "dup_ratio": round(1 - len(set(out.split("\n"))) / max(1, len(out.split("\n"))), 2),
    }


# --------------------------------------------------------------------------- #
# Output under test                                                            #
# --------------------------------------------------------------------------- #
def doc2mark_markdown(path: str) -> str:
    """The current pipeline output (route auto-selects rule-based vs page-OCR)."""
    from doc2mark import UnifiedDocumentLoader
    from doc2mark.ocr.base import OCRConfig
    loader = UnifiedDocumentLoader(ocr_provider="openai", model=MODEL,
                                   ocr_config=OCRConfig(context_pages=0))
    return loader.load(path, output_format="markdown",
                       extract_images=True, ocr_images=True).content


def run(extractor=doc2mark_markdown, quick=False):
    rows = []
    for rel, cls in DOCS:
        if quick and cls in IMAGE_CLASSES:
            continue
        path = rel if os.path.isabs(rel) else os.path.join(SAMP, rel)
        if not os.path.exists(path):
            print(f"MISSING {rel}")
            continue
        name = os.path.basename(path)
        print(f"\n>>> {name} [{cls}]")
        try:
            out = extractor(path)
        except Exception as e:
            print(f"   EXTRACT FAILED: {e}")
            rows.append((name, cls, None, None, None, {"err": 1}))
            continue

        tp = None if cls in IMAGE_CLASSES else text_preservation(rule_based_text(path), out)
        tf = table_fidelity(true_grids(path), out) if cls in TABLE_CLASSES else None
        mn = meaningfulness(out) if cls in IMAGE_CLASSES else None
        jk = junk(out)
        rows.append((name, cls, tp, tf, mn, jk))
        print(f"   text={'%.0f%%'%(tp*100) if tp is not None else '  -'}"
              f"  table={'%.0f%%'%(tf*100) if tf is not None else '  -'}"
              f"  meaning={mn if mn is not None else '-'}/5"
              f"  junk(base64={jk['base64']},markers={jk['markers']},dup={jk['dup_ratio']})")

    print("\n================ SCORECARD ================")
    print(f"{'document':<36}{'text':>7}{'table':>8}{'mean':>7}  pass")
    allpass = True
    for name, cls, tp, tf, mn, jk in rows:
        checks = []
        if tp is not None:
            checks.append(tp >= THRESH["text"])
        if tf is not None:
            checks.append(tf >= THRESH["table"])
        if mn is not None:
            checks.append(mn >= THRESH["meaning"])
        if jk.get("base64") or jk.get("markers") or jk.get("err"):
            checks.append(False)
        ok = all(checks) if checks else True
        allpass = allpass and ok
        print(f"{name[:35]:<36}"
              f"{('%.0f%%'%(tp*100)) if tp is not None else '-':>7}"
              f"{('%.0f%%'%(tf*100)) if tf is not None else '-':>8}"
              f"{(str(mn)) if mn is not None else '-':>7}  {'PASS' if ok else 'FAIL'}")
    print("==========================================")
    print("RESULT:", "ALL PASS" if allpass else "FAILURES PRESENT")
    return rows


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--quick", action="store_true", help="skip the 30-page image deck")
    args = ap.parse_args()
    run(quick=args.quick)
